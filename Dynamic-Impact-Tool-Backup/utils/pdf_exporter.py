from fpdf import FPDF
import tempfile
import os

class PDF(FPDF):
    def header(self):
        self.set_font("Arial", "B", 14)
        self.cell(0, 10, "Dynamic Impact Tool Report", ln=True, align="C")
        self.ln(5)

    def chapter_title(self, title):
        self.set_font("Arial", "B", 12)
        self.set_text_color(0, 0, 128)
        self.cell(0, 10, title, ln=True, align="L")
        self.ln(2)

    def chapter_body(self, text):
        self.set_font("Arial", "", 11)
        self.set_text_color(0)
        self.multi_cell(0, 8, text)
        self.ln()

def generate_pdf_report(insights, chat_logs):
    pdf = PDF()
    pdf.add_page()

    if insights:
        pdf.chapter_title("📝 Dataset Insight Summary")
        pdf.chapter_body(insights)

    if chat_logs:
        pdf.chapter_title("💬 Chat History")
        for item in chat_logs:
            user = item.get("user", "")
            assistant = item.get("assistant", {}).get("response", item.get("assistant", ""))
            pdf.set_font("Arial", "B", 11)
            pdf.multi_cell(0, 8, f"User: {user}")
            pdf.set_font("Arial", "", 11)
            pdf.multi_cell(0, 8, f"AI: {assistant}")
            pdf.ln(3)

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
            pdf.output(tmp_file.name)
            return tmp_file.name
    except Exception as e:
        print(f"[PDF ERROR] {e}")
        return None
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.shapes import MSO_SHAPE
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# import tempfile
# import os

# def generate_pdf_report(insights: str, chat_logs: list, filename="summary.pdf"):
#     file_path = os.path.join(tempfile.gettempdir(), filename)
#     c = canvas.Canvas(file_path, pagesize=letter)
#     width, height = letter

#     c.setFont("Helvetica-Bold", 16)
#     c.drawString(40, height - 50, "🧠 Insight Summary")
#     c.setFont("Helvetica", 12)

#     text_obj = c.beginText(40, height - 80)
#     for line in insights.split('\n'):
#         text_obj.textLine(line)
#     c.drawText(text_obj)

#     c.setFont("Helvetica-Bold", 16)
#     c.drawString(40, height - 280, "💬 Chat Logs")
#     text_obj = c.beginText(40, height - 310)
#     for chat in chat_logs:
#         text_obj.textLine(f"You: {chat['user']}")
#         response = chat['assistant'].get("response", str(chat['assistant']))
#         for line in response.split('\n'):
#             text_obj.textLine(f"AI: {line}")
#         text_obj.textLine("")

#     c.drawText(text_obj)
#     c.showPage()
#     c.save()
#     return file_path


# def export_to_pptx(insights: str, chat_logs: list, filename="summary.pptx"):
#     prs = Presentation()
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])
#     title_slide.shapes.title.text = "Dynamic Impact Tool Summary"
#     title_slide.placeholders[1].text = "Insights & Chat Logs"

#     # Slide for insights
#     slide = prs.slides.add_slide(prs.slide_layouts[1])
#     slide.shapes.title.text = "🧠 Insights"
#     content = slide.placeholders[1]
#     content.text = insights if insights else "No insights available."

#     # Slide for chats
#     chat_slide = prs.slides.add_slide(prs.slide_layouts[1])
#     chat_slide.shapes.title.text = "💬 Chat Logs"
#     chat_text = ""
#     for chat in chat_logs:
#         user = chat["user"]
#         assistant = chat["assistant"].get("response", str(chat["assistant"]))
#         chat_text += f"You: {user}\nAI: {assistant}\n\n"
#     chat_slide.placeholders[1].text = chat_text[:4000]  # Limit text due to PowerPoint constraints

#     file_path = os.path.join(tempfile.gettempdir(), filename)
#     prs.save(file_path)
#     return file_path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import tempfile
import os

def generate_pdf_report(insights=None, chat_history=None, comparison_history=None, charts=None, metadata=None, filename="summary.pdf"):
    file_path = os.path.join(tempfile.gettempdir(), filename)
    c = canvas.Canvas(file_path, pagesize=letter)
    width, height = letter

    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, height - 50, "🧠 Insight Summary")
    c.setFont("Helvetica", 12)

    text_obj = c.beginText(40, height - 80)
    for line in (insights or "No insights available.").split('\n'):
        text_obj.textLine(line)
    c.drawText(text_obj)

    if chat_history:
        c.setFont("Helvetica-Bold", 16)
        c.drawString(40, height - 280, "💬 Chat Logs")
        text_obj = c.beginText(40, height - 310)
        for chat in chat_history:
            text_obj.textLine(f"You: {chat['user']}")
            response = chat['assistant'].get("response", str(chat['assistant']))
            for line in response.split('\n'):
                text_obj.textLine(f"AI: {line}")
            text_obj.textLine("")
        c.drawText(text_obj)

    c.showPage()
    c.save()
    return file_path

def export_to_pptx(insights=None, chat_history=None, comparison_history=None, charts=None, metadata=None, filename="summary.pptx"):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Dynamic Impact Tool Summary"
    title_slide.placeholders[1].text = "Insights & Chat Logs"

    # Slide for insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "🧠 Insights"
    content = slide.placeholders[1]
    content.text = insights if insights else "No insights available."

    # Slide for chats
    if chat_history:
        chat_slide = prs.slides.add_slide(prs.slide_layouts[1])
        chat_slide.shapes.title.text = "💬 Chat Logs"
        chat_text = ""
        for chat in chat_history:
            user = chat["user"]
            assistant = chat["assistant"].get("response", str(chat["assistant"]))
            chat_text += f"You: {user}\nAI: {assistant}\n\n"
        chat_slide.placeholders[1].text = chat_text[:4000]  # Limit text due to PowerPoint constraints

    file_path = os.path.join(tempfile.gettempdir(), filename)
    prs.save(file_path)
    return file_path
