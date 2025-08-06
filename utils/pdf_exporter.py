import os
import plotly.express as px
from tempfile import NamedTemporaryFile
from fpdf import FPDF
from fpdf.enums import XPos, YPos
import os
from datetime import datetime
from utils.sections.summary import generate_section1_summary
from utils.sections.introduction import generate_section2_introduction
from utils.sections.data_overview import generate_section3_data_overview
from utils.sections.methodology import generate_section4_methodology
from utils.sections.analysis_insights import generate_section5_analysis_insights
from utils.sections.cross_domain import generate_section6_cross_domain
from utils.sections.recommendation import generate_section7_recommendations
from utils.sections.conclusion import generate_section8_conclusion
import textwrap
import re
import streamlit as st

class CustomPDF(FPDF):
    def footer(self):
        if self.page_no() > 1:
            self.set_y(-15)
            self.set_font("DejaVu", "I", 10)
            self.cell(0, 10, f"Page {self.page_no()}", 0, 0, "C")

import re

def clean_insight_text(text):
    text = re.sub(
        r"(\n\|.+?\|\n\|[-| :]+\|\n(?:\|.*?\|\n?)+)",
        "", text, flags=re.MULTILINE
    )
    lines = text.split("\n")
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped or "Analytical Insight" in stripped or re.match(r"=+", stripped):
            continue
        cleaned_lines.append(stripped)
    return "\n".join(cleaned_lines)

def safe_multicell(pdf, text, width=180, font_size=12):
    import textwrap
    import re

    def clean(text):
        return str(text).encode("utf-8", errors="ignore").decode("utf-8")

    pdf.set_font("DejaVu", "", font_size)

    text = clean(text)
    text = re.sub(r"(?i)^analytical insight:.*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"={3,}|-{3,}", "", text)
    text = re.sub(r"\*\*+", "", text)
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    text = text.replace("🚀", "").replace("🧑", "").replace("🤖", "")

    paragraphs = text.split("\n")
    for para in paragraphs:
        para = para.strip()
        if para:
            wrapped_lines = textwrap.wrap(para, width=90)
            for line in wrapped_lines:
                pdf.cell(0, 10, clean(line), ln=True)
            pdf.ln(2)


import re

def extract_markdown_tables(text):
    pattern = r"((\|.+?\|)+\n(\|[-:]+?\|)+\n((\|.*?\|)+\n?)+)"
    return re.findall(pattern, text)

def render_markdown_table(pdf, table_text):
    lines = [line.strip() for line in table_text.strip().split('\n') if line.strip()]
    if not lines or len(lines) < 2:
        return

    headers = [cell.strip() for cell in lines[0].split('|') if cell.strip()]
    rows = [
        [cell.strip() for cell in line.split('|') if cell.strip()]
        for line in lines[2:]
    ]

    col_width = pdf.w / (len(headers) + 1)  # Adjust width

    # Render header
    pdf.set_font("DejaVu", "B", 11)
    for header in headers:
        pdf.cell(col_width, 10, header, border=1)
    pdf.ln()

    # Render rows
    pdf.set_font("DejaVu", "", 11)
    for row in rows:
        for cell in row:
            pdf.cell(col_width, 10, cell, border=1)
        pdf.ln()

    pdf.ln(5)


def draw_markdown_table(pdf, markdown_text):
    lines = markdown_text.strip().split("\n")
    table_lines = [line for line in lines if "|" in line and "---" not in line]

    if not table_lines:
        return False  # No table detected

    # Clean and split table rows
    rows = [line.strip().strip("|").split("|") for line in table_lines]
    col_width = 180 // len(rows[0])
    pdf.set_font("DejaVu", "", 11)
    

    for row in rows:
        for cell in row:
            pdf.cell(col_width, 10, cell.strip(), border=1)
        pdf.ln()
    pdf.ln(5)
    return True

def generate_pdf_report(session, filename="summary.pdf", model_source="groq"):
    df = session["df"]
    dataset_name = session.get("name", "Unnamed Dataset")
    insights_raw = session.get("selected_insight_results", [])
    insights = list(insights_raw.values()) if isinstance(insights_raw, dict) else insights_raw
    chat_history = session.get("chat_history", [])

    # Set up PDF
    pdf = CustomPDF()

    font_base_path = os.path.join("assets", "fonts", "dejavu-fonts-ttf-2.37", "ttf")
    font_path = os.path.join(font_base_path, "DejaVuSans.ttf")
    font_bold_path = os.path.join(font_base_path, "DejaVuSans-Bold.ttf")
    font_italic_path = os.path.join(font_base_path, "DejaVuSans-Oblique.ttf")

    pdf.add_font("DejaVu", "", font_path, uni=True)
    pdf.add_font("DejaVu", "B", font_bold_path, uni=True)
    pdf.add_font("DejaVu", "I", font_italic_path, uni=True)

    pdf.set_font("DejaVu", "", 14)
    pdf.set_auto_page_break(auto=True, margin=15)

    section_pages = []  # to store (section title, page no)

    # === COVER PAGE ===
    pdf.add_page()
    pdf.set_font("DejaVu", "B", 20)
    pdf.cell(0, 20, "Dynamic Dataset Analysis Report", ln=True, align="C")
    pdf.set_font("DejaVu", "", 16)
    pdf.cell(0, 10, "Generated by Dynamic Impact Tool", ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("DejaVu", "", 12)
    pdf.cell(0, 10, f"Date: {datetime.now().strftime('%d %B %Y')}", ln=True, align="C")
    pdf.cell(0, 10, f"Dataset: {dataset_name}", ln=True, align="C")

    # === INDEX PAGE ===
    pdf.add_page()
    section_pages.append(("Index", pdf.page_no()))

    try:
        summary_text = generate_section1_summary(df, model_source)
        pdf.add_page()
        section_pages.append(("Executive Summary", pdf.page_no()))
        pdf.set_font("DejaVu", "B", 16)
        pdf.cell(0, 10, "1. Executive Summary", ln=True)
        pdf.set_font("DejaVu", "", 12)
        safe_multicell(pdf, summary_text)
    except Exception as e:
        st.error(f"PDF export error in Section 1: {e}")

    # === SECTION 2 ===
    pdf.add_page()
    section_pages.append(("Introduction", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "2. Introduction", ln=True)
    safe_multicell(pdf, generate_section2_introduction(df, model_source))

    # === SECTION 3 ===
    pdf.add_page()
    section_pages.append(("Data Overview", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "3. Data Overview", ln=True)
    safe_multicell(pdf, generate_section3_data_overview(df, model_source))

    # === SECTION 4 ===
    pdf.add_page()
    section_pages.append(("Methodology", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "4. Methodology", ln=True)
    safe_multicell(pdf, generate_section4_methodology(df, model_source))

        # === SECTION 5: Detailed Analysis & Insights ===
    try:
        analysis_text = generate_section5_analysis_insights(df, model_source)
        
        # Add new page and section heading
        pdf.add_page()
        section_pages.append(("Detailed Analysis & Insights", pdf.page_no()))
        pdf.set_font("DejaVu", "B", 16)
        pdf.cell(0, 10, "5. Detailed Analysis & Insights", ln=True)
        pdf.set_font("DejaVu", "", 12)

        # Add AI-generated text
        if analysis_text:
            safe_multicell(pdf, analysis_text)
        else:
            safe_multicell(pdf, "No insights were generated for this dataset.")

        if 'year' in df.columns and 'number_of_customers_per_day' in df.columns:
            df_yearly = df.groupby('year')['number_of_customers_per_day'].mean().reset_index()
            fig = px.bar(
                df_yearly,
                x='year',
                y='number_of_customers_per_day',
                title='Average Number of Customers per Day by Year',
                labels={'number_of_customers_per_day': 'Avg Customers/Day'}
            )

            # Save chart to temporary image file
            with NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
                fig.write_image(tmpfile.name, width=800, height=500)
                chart_path = tmpfile.name

            # Add chart image to PDF
            pdf.ln(5)
            pdf.image(chart_path, x=None, w=pdf.w - 30)  # Auto center
            os.remove(chart_path)  # Clean up after adding

    except Exception as e:
        st.error(f"PDF export error in Section 5 (Detailed Analysis & Insights): {e}")

    try:
     cross_domain_text = generate_section6_cross_domain(df, model_source)
     pdf.add_page()
     section_pages.append(("Cross-Domain Insights", pdf.page_no()))
     pdf.set_font("DejaVu", "B", 16)
     pdf.cell(0, 10, "6. Cross-Domain Insights", ln=True)
     pdf.set_font("DejaVu", "", 12)
     safe_multicell(pdf, cross_domain_text)
    except Exception as e:
    # Check for Groq API token limit error and skip if matched
     if "Groq API error" in str(e) and "rate_limit_exceeded" in str(e):
        st.warning("Section 6 (Cross-Domain Insights) skipped due to Groq API token limit.")
        # Skip adding to section_pages (so it's not added to index)
     else:
        # Show error for debugging
        st.error(f"PDF export error in Section 6: {e}")
    # === SECTION 7: Recommendations & Actionable Items ===
    recommendations = generate_section7_recommendations(insights, model_source)
    pdf.add_page()
    section_pages.append(("Recommendations & Actionable Items", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "7. Recommendations & Actionable Items", ln=True)
    pdf.ln(5)

    if recommendations:
     pdf.set_font("DejaVu", "", 12)
     for idx, rec in enumerate(recommendations, start=1):
        if isinstance(rec, dict):
            text = (
                f"Insight: {rec.get('Insight', 'N/A')}\n"
                f"Recommended Action: {rec.get('Recommended Action', 'N/A')}\n"
                f"Priority: {rec.get('Priority', 'N/A')}\n"
                f"Owner/Team: {rec.get('Owner', 'N/A')}\n"
                f"Timeline: {rec.get('Timeline', 'N/A')}\n"
                "----------------------------------------------------\n"
            )
        else:
            text = f"{str(rec)}\n----------------------------------------------------\n"
        safe_multicell(pdf, text)
    else:
     pdf.set_font("DejaVu", "", 12)
     pdf.multi_cell(0, 10, "No recommendations were generated.")
    # === SECTION 8: Conclusion ===
    try:
        conclusion_text = generate_section8_conclusion(
            df,
            insights=session.get("selected_insight_results", {}),
            model_source=model_source
        )
        pdf.add_page()
        section_pages.append(("Conclusion", pdf.page_no()))
        pdf.set_font("DejaVu", "B", 16)
        pdf.cell(0, 10, "8. Conclusion", ln=True)
        pdf.set_font("DejaVu", "", 12)
        safe_multicell(pdf, conclusion_text)
    except Exception as e:
        st.error(f"PDF export error in Section 8 (Conclusion): {e}")



    # === INDEX PAGE REWRITE ===
    pdf.page = 2  # Set to index page (2nd page)
    pdf.set_xy(10, 20)
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "Index", ln=True, align="C")
    pdf.ln(5)
    pdf.set_font("DejaVu", "B", 12)
    pdf.cell(20, 10, "S.No", border=1)
    pdf.cell(110, 10, "Section Title", border=1)
    pdf.cell(30, 10, "Page No.", border=1, ln=True)
    pdf.set_font("DejaVu", "", 12)
    for i, (title, page) in enumerate(section_pages[1:], start=1):  # skip index
        pdf.cell(20, 10, str(i), border=1)
        pdf.cell(110, 10, title, border=1)
        pdf.cell(30, 10, str(page), border=1, ln=True)

    # === Export ===
    os.makedirs("exports", exist_ok=True)
    output_path = os.path.join("exports", filename)
    pdf.output(output_path)
    return output_path


# ==================== PPTX EXPORTER ==================== #
from pptx import Presentation
def export_to_pptx(session, filename="summary.pptx"):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"Dataset: {session.get('name', 'Unnamed Dataset')}"
    insight_list = session.get("insights") or session.get("selected_insight_results", [])
    for insight in insight_list:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{insight['question']}"
        slide.shapes.placeholders[1].text = insight['result']

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Chat History"
    chat_text = ""
    for chat in session.get('chat_history', []):
        chat_text += f"🧑 {chat['user']}\n🤖 {chat['assistant'].get('response', '')}\n\n"

    slide.shapes.placeholders[1].text = chat_text

    output_path = os.path.join("exports", filename)
    prs.save(output_path)
    return output_path



